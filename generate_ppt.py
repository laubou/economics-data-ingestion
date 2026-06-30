"""
Generates the IATA Economics Pipeline presentation.
Run: python generate_ppt.py
Output: IATA_Economics_Pipeline.pptx
"""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ─── Palette ────────────────────────────────────────────────────────────────

NAVY        = RGBColor(0x00, 0x2B, 0x5C)  # IATA dark blue
BLUE        = RGBColor(0x00, 0x6C, 0xA9)  # mid blue
SKY         = RGBColor(0x00, 0xAA, 0xD4)  # light blue
ORANGE      = RGBColor(0xE8, 0x6A, 0x1A)  # accent / highlight
GREEN       = RGBColor(0x27, 0x9A, 0x5E)  # success / check
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xF4, 0xF7, 0xFA)
DARK_GRAY   = RGBColor(0x2D, 0x32, 0x3B)
MID_GRAY    = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY  = RGBColor(0xE2, 0xE8, 0xF0)

W = 13.33   # slide width  (inches)
H = 7.5     # slide height (inches)

# ─── Presentation bootstrap ─────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = Inches(W)
prs.slide_height = Inches(H)
blank_layout = prs.slide_layouts[6]   # truly blank


# ─── Low-level helpers ───────────────────────────────────────────────────────

def rgb(r: int, g: int, b: int) -> RGBColor:
    return RGBColor(r, g, b)


def rect(slide, x, y, w, h, fill=WHITE, border=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def text_box(slide, text, x, y, w, h, size=18, bold=False, color=DARK_GRAY,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb


def label_in_rect(slide, text, x, y, w, h, fill=BLUE, text_color=WHITE,
                  size=13, bold=True, align=PP_ALIGN.CENTER, radius=None):
    """Colored rectangle with centered text on top."""
    shape = rect(slide, x, y, w, h, fill=fill)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    run.font.name = "Calibri"
    return shape


def arrow(slide, x1, y1, x2, y2, color=MID_GRAY):
    """Simple horizontal or vertical connector arrow."""
    from pptx.util import Emu
    from pptx.oxml.ns import qn
    import lxml.etree as etree

    cx = abs(Inches(x2) - Inches(x1))
    cy = abs(Inches(y2) - Inches(y1))
    left  = Inches(min(x1, x2))
    top   = Inches(min(y1, y2))

    connector = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = color
    connector.line.width = Pt(1.5)
    return connector


def header_band(slide, title, subtitle=None):
    """Dark navy header bar at top of slide."""
    rect(slide, 0, 0, W, 1.1, fill=NAVY)
    text_box(slide, title, 0.4, 0.1, W - 0.8, 0.6,
             size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        text_box(slide, subtitle, 0.4, 0.65, W - 0.8, 0.38,
                 size=14, bold=False, color=SKY, align=PP_ALIGN.LEFT)
    # thin accent line
    rect(slide, 0, 1.1, W, 0.04, fill=ORANGE)


def bullet_block(slide, items, x, y, w, size=16, color=DARK_GRAY, spacing=0.38):
    """Stack of bullet lines."""
    for i, item in enumerate(items):
        text_box(slide, item, x, y + i * spacing, w, 0.35, size=size, color=color)


def tag(slide, text, x, y, fill=SKY, text_color=WHITE, size=11):
    """Small pill-shaped label."""
    label_in_rect(slide, text, x, y, 1.4, 0.3, fill=fill,
                  text_color=text_color, size=size, bold=False)


def add_notes(slide, notes_text: str):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


# ═══════════════════════════════════════════════════════════════════════════ #
# SLIDE 1 — Title                                                              #
# ═══════════════════════════════════════════════════════════════════════════ #

s = prs.slides.add_slide(blank_layout)

# Full background
rect(s, 0, 0, W, H, fill=NAVY)
# Bottom accent stripe
rect(s, 0, H - 0.5, W, 0.5, fill=ORANGE)
# White card
rect(s, 0.7, 1.2, W - 1.4, 4.8, fill=WHITE)

text_box(s, "IATA Economics", 0.9, 1.4, W - 1.8, 1.0,
         size=40, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
text_box(s, "Data Ingestion Pipeline", 0.9, 2.3, W - 1.8, 0.9,
         size=36, bold=False, color=BLUE, align=PP_ALIGN.LEFT)

rect(s, 0.9, 3.25, 4.5, 0.05, fill=ORANGE)

text_box(s, "Case Study — Production-Ready Streaming Architecture on AWS",
         0.9, 3.5, W - 1.8, 0.6, size=16, color=MID_GRAY, align=PP_ALIGN.LEFT)
text_box(s, "Laura Bourbon  ·  2026", 0.9, 4.9, W - 1.8, 0.4,
         size=13, color=MID_GRAY, align=PP_ALIGN.LEFT)

add_notes(s, "Bonjour, je vais vous présenter mon case study sur la construction d'un pipeline d'ingestion de données pour l'équipe Economics de l'IATA. Le pipeline est production-ready, déployable sur AWS via Terraform, et conçu pour évoluer d'un fichier one-off vers une ingestion continue quotidienne.")


# ═══════════════════════════════════════════════════════════════════════════ #
# SLIDE 2 — Agenda                                                             #
# ═══════════════════════════════════════════════════════════════════════════ #

s = prs.slides.add_slide(blank_layout)
rect(s, 0, 0, W, H, fill=OFF_WHITE)
header_band(s, "Agenda", "15 minutes · 4 sections · Q&A")

items = [
    ("01", "Business Context & Objective",    "1 min"),
    ("02", "Architecture Overview",           "2 min"),
    ("03", "Deep Dive — 4-Step Dataflow",     "7 min"),
    ("04", "Infrastructure as Code",          "2 min"),
    ("05", "Key Design Decisions",            "2 min"),
    ("  ", "Q&A",                             "1 min"),
]

for i, (num, title, dur) in enumerate(items):
    y = 1.4 + i * 0.82
    fill = NAVY if i < 5 else MID_GRAY
    label_in_rect(s, num, 0.5, y, 0.55, 0.55, fill=fill, size=18, bold=True)
    text_box(s, title, 1.25, y + 0.05, 8.5, 0.5, size=18, bold=(i < 5), color=DARK_GRAY)
    text_box(s, dur,   10.5, y + 0.05, 2.0, 0.5, size=14, color=MID_GRAY, align=PP_ALIGN.RIGHT)

add_notes(s, "On va couvrir : le contexte business en 1 minute, l'architecture globale, puis un deep dive sur chaque étape du dataflow. On termine sur l'IaC et les choix de conception. Laissez vos questions pour la fin.")


# ═══════════════════════════════════════════════════════════════════════════ #
# SLIDE 3 — Business Context                                                   #
# ═══════════════════════════════════════════════════════════════════════════ #

s = prs.slides.add_slide(blank_layout)
rect(s, 0, 0, W, H, fill=OFF_WHITE)
header_band(s, "Business Context", "What we need to build and why")

# Left column - situation
label_in_rect(s, "SITUATION", 0.4, 1.3, 5.8, 0.4, fill=BLUE, size=13)
rect(s, 0.4, 1.7, 5.8, 3.5, fill=WHITE, border=LIGHT_GRAY)

situation = [
    "📦  Source: CSV file with 2M+ sales records",
    "      Provider URL (zipped) — updated daily",
    "",
    "📋  Fields: 14 business columns",
    "      Orders, countries, revenues, margins…",
    "",
    "🎯  Target: Queryable data lake on AWS",
    "      Bronze (raw) + Silver (curated)",
    "",
    "🔄  Continuous: daily batch from 2026",
]
for i, line in enumerate(situation):
    text_box(s, line, 0.6, 1.85 + i * 0.31, 5.4, 0.32, size=13, color=DARK_GRAY)

# Right column - requirements
label_in_rect(s, "5 REQUIREMENTS", 6.5, 1.3, 6.4, 0.4, fill=ORANGE, size=13)
rect(s, 6.5, 1.7, 6.4, 3.5, fill=WHITE, border=LIGHT_GRAY)

reqs = [
    ("✅", "Acquire & Land", "Batch trigger, archive + land on S3"),
    ("✅", "Stream into Bronze", "Real-time, Iceberg table, Glue catalog"),
    ("✅", "Transform → Silver", "Typed, deduped, partitioned"),
    ("✅", "Serve", "Amazon Athena queries on bronze & silver"),
    ("✅", "IaC", "All resources deployable via Terraform"),
]
for i, (icon, title, sub) in enumerate(reqs):
    y = 1.85 + i * 0.68
    text_box(s, icon, 6.7, y, 0.4, 0.35, size=18)
    text_box(s, title, 7.2, y, 5.2, 0.3, size=14, bold=True, color=NAVY)
    text_box(s, sub, 7.2, y + 0.3, 5.5, 0.3, size=12, color=MID_GRAY)

add_notes(s, "L'équipe Economics a un fichier de 2M records accessible via une URL. L'objectif est de construire un pipeline complet : téléchargement, archivage, streaming, transformation et service via Athena. Et la contrainte clé : le fournisseur va commencer à envoyer des données en continu — donc ce n'est pas un script jetable, c'est le point de départ d'un vrai pipeline.")


# ═══════════════════════════════════════════════════════════════════════════ #
# SLIDE 4 — Architecture Overview                                              #
# ═══════════════════════════════════════════════════════════════════════════ #

s = prs.slides.add_slide(blank_layout)
rect(s, 0, 0, W, H, fill=OFF_WHITE)
header_band(s, "Architecture Overview", "Medallion architecture on AWS — landing → bronze → silver")

# ── Row labels
def row_label(slide, text, y):
    rect(slide, 0.15, y, 1.1, 0.42, fill=NAVY)
    text_box(slide, text, 0.15, y, 1.1, 0.42, size=10, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)

row_label(s, "TRIGGER", 1.25)
row_label(s, "BATCH", 2.0)
row_label(s, "STREAM", 2.85)
row_label(s, "LAKE", 3.85)
row_label(s, "SERVE", 5.05)

# ── Trigger row
label_in_rect(s, "EventBridge\nScheduler\n09:00 UTC", 1.4, 1.2, 2.0, 0.62,
              fill=SKY, size=11, bold=True)
arrow(s, 3.4, 1.5, 4.0, 1.5)
label_in_rect(s, "Step Functions\nState Machine", 4.0, 1.2, 2.0, 0.62,
              fill=SKY, size=11, bold=False)
label_in_rect(s, "IngestionState\n(S3 ETag tracking)", 9.0, 1.2, 2.2, 0.62,
              fill=LIGHT_GRAY, text_color=DARK_GRAY, size=11, bold=False)

# ── Batch row
label_in_rect(s, "ECS Fargate\nDownloader", 1.4, 2.0, 2.1, 0.55,
              fill=BLUE, size=12)
arrow(s, 3.5, 2.25, 4.1, 2.25)
label_in_rect(s, "S3\narchive/  landing/", 4.1, 2.0, 2.4, 0.55,
              fill=rgb(0xE9, 0x6B, 0x1A), size=12)
arrow(s, 6.5, 2.25, 7.1, 2.25)
label_in_rect(s, "ECS Fargate\nProducer", 7.1, 2.0, 2.1, 0.55,
              fill=BLUE, size=12)
arrow(s, 9.2, 2.25, 9.8, 2.25)
label_in_rect(s, "Amazon MSK\n(Kafka)", 9.8, 2.0, 1.9, 0.55,
              fill=NAVY, size=12)

# ── Stream row
label_in_rect(s, "Amazon MSK\n(Kafka)", 1.4, 2.85, 1.9, 0.55,
              fill=NAVY, size=12)
arrow(s, 3.3, 3.1, 3.9, 3.1)
label_in_rect(s, "ECS Fargate\nBronze Consumer", 3.9, 2.85, 2.4, 0.55,
              fill=BLUE, size=12)
arrow(s, 6.3, 3.1, 6.9, 3.1)
label_in_rect(s, "S3\nbronze/ (Iceberg)", 6.9, 2.85, 2.3, 0.55,
              fill=rgb(0xE9, 0x6B, 0x1A), size=12)
arrow(s, 9.2, 3.1, 9.8, 3.1)
label_in_rect(s, "Glue Data\nCatalog", 9.8, 2.85, 1.9, 0.55,
              fill=GREEN, size=12)

# ── Lake row
label_in_rect(s, "S3 bronze/\n(Iceberg Parquet)", 1.4, 3.85, 2.4, 0.6,
              fill=rgb(0xE9, 0x6B, 0x1A), size=12)
arrow(s, 3.8, 4.1, 4.4, 4.1)
label_in_rect(s, "ECS Fargate\nTransformer Silver", 4.4, 3.85, 2.4, 0.6,
              fill=BLUE, size=12)
arrow(s, 6.8, 4.1, 7.4, 4.1)
label_in_rect(s, "S3 silver/\n(Iceberg Parquet)", 7.4, 3.85, 2.4, 0.6,
              fill=rgb(0xE9, 0x6B, 0x1A), size=12)
arrow(s, 9.8, 4.1, 10.4, 4.1)
label_in_rect(s, "Glue Data\nCatalog", 10.4, 3.85, 1.9, 0.6,
              fill=GREEN, size=12)

# ── Serve row
label_in_rect(s, "Glue Data Catalog\n(sales_bronze + sales_silver)", 1.4, 5.05, 3.2, 0.6,
              fill=GREEN, size=12)
arrow(s, 4.6, 5.35, 5.2, 5.35)
label_in_rect(s, "Amazon Athena\nEngine v3 Workgroup", 5.2, 5.05, 3.0, 0.6,
              fill=SKY, size=12)
arrow(s, 8.2, 5.35, 8.8, 5.35)
label_in_rect(s, "Named Queries\n+ Cost Control", 8.8, 5.05, 2.4, 0.6,
              fill=MID_GRAY, size=12)

# Legend
rect(s, 0.15, 6.6, W - 0.3, 0.55, fill=WHITE, border=LIGHT_GRAY)
text_box(s, "■ ECS Fargate Services", 0.4, 6.67, 2.5, 0.38, size=11, color=BLUE)
text_box(s, "■ S3 Data Lake",         3.0, 6.67, 2.0, 0.38, size=11, color=ORANGE)
text_box(s, "■ AWS Managed Services", 5.0, 6.67, 2.5, 0.38, size=11, color=GREEN)
text_box(s, "■ Orchestration",        7.5, 6.67, 2.0, 0.38, size=11, color=SKY)
text_box(s, "Python shared package: economics_pipeline", 9.5, 6.67, 3.5, 0.38, size=11, color=MID_GRAY)

add_notes(s, "Voici l'architecture complète. On voit 5 couches : le trigger (EventBridge), le batch (downloader + producer), le streaming (MSK + consumer), le lac de données avec les tables Iceberg, et la couche de service Athena. Tout est découplé — si le volume explose, on scale chaque couche indépendamment. Le code Python est dans un package partagé installé dans chaque conteneur ECS.")


# ═══════════════════════════════════════════════════════════════════════════ #
# SLIDE 5 — Step 1: Acquire & Land                                             #
# ═══════════════════════════════════════════════════════════════════════════ #

s = prs.slides.add_slide(blank_layout)
rect(s, 0, 0, W, H, fill=OFF_WHITE)
header_band(s, "Step 1 — Acquire & Land", "Batch trigger · Idempotent · ETag-based change detection")

# Left: mini flow
label_in_rect(s, "EventBridge Scheduler\ncron(0 9 * * ? *)", 0.4, 1.3, 3.2, 0.7,
              fill=SKY, size=13)
arrow(s, 2.0, 2.0, 2.0, 2.4)
label_in_rect(s, "ECS Fargate\nDownloader", 0.4, 2.4, 3.2, 0.65, fill=BLUE, size=13)
arrow(s, 1.0, 3.05, 1.0, 3.4)
arrow(s, 2.8, 3.05, 2.8, 3.4)
label_in_rect(s, "S3  archive/dataset.zip", 0.4, 3.4, 1.9, 0.6,
              fill=ORANGE, size=12, bold=False)
label_in_rect(s, "S3  landing/dataset.csv", 2.5, 3.4, 1.9, 0.6,
              fill=ORANGE, size=12, bold=False)
arrow(s, 2.0, 4.0, 2.0, 4.35)
label_in_rect(s, "IngestionState  ✓  saved", 0.4, 4.35, 3.2, 0.55,
              fill=GREEN, size=13)

# Right: key points
label_in_rect(s, "KEY POINTS", 4.0, 1.3, 8.9, 0.4, fill=NAVY, size=13)

points = [
    ("ETag detection",
     "HTTP HEAD before every download. If ETag matches last run → skip entirely.\n"
     "Prevents redundant downloads when provider hasn't updated the file."),
    ("Idempotent state",
     "IngestionStateManager persists {url, etag, checksum_md5, status} to S3.\n"
     "Status: 'extracted' → 'produced'. Interrupted run resumes safely."),
    ("archive/ vs landing/",
     "archive/ = original zip, never modified (audit trail).\n"
     "landing/ = extracted CSV, ready for the producer."),
    ("Retry logic",
     "@network_retry: 3 attempts, exponential 1→8s backoff.\n"
     "Raises DownloadError with original cause on final failure."),
]
for i, (title, body) in enumerate(points):
    y = 1.9 + i * 1.22
    rect(s, 4.0, y, 8.9, 1.1, fill=WHITE, border=LIGHT_GRAY)
    text_box(s, title, 4.2, y + 0.07, 8.5, 0.33, size=14, bold=True, color=NAVY)
    text_box(s, body,  4.2, y + 0.38, 8.5, 0.65, size=12, color=DARK_GRAY)

add_notes(s, "Le downloader commence toujours par un HTTP HEAD sur l'URL source. Il compare l'ETag retourné avec celui sauvegardé dans l'IngestionStateManager. Si c'est identique, on sort immédiatement — aucun re-téléchargement inutile. Si c'est nouveau, on télécharge, on archive le zip, on extrait le CSV, et on sauvegarde le nouvel état avec le statut 'extracted'. Le producer verra ce statut 'extracted' et sait qu'il y a du nouveau à produire.")


# ═══════════════════════════════════════════════════════════════════════════ #
# SLIDE 6 — Step 2: Stream into Bronze                                         #
# ═══════════════════════════════════════════════════════════════════════════ #

s = prs.slides.add_slide(blank_layout)
rect(s, 0, 0, W, H, fill=OFF_WHITE)
header_band(s, "Step 2 — Stream into Bronze", "CSV → Kafka/MSK → Iceberg table registered in Glue")

# Left: flow
label_in_rect(s, "S3 landing/\ndataset.csv", 0.4, 1.3, 2.6, 0.65,
              fill=ORANGE, size=13)
arrow(s, 1.7, 1.95, 1.7, 2.35)
label_in_rect(s, "Producer (ECS)\ncsv.DictReader row by row", 0.4, 2.35, 2.6, 0.65,
              fill=BLUE, size=12)
arrow(s, 1.7, 3.0, 1.7, 3.35)
label_in_rect(s, "Amazon MSK\nkey = order_id", 0.4, 3.35, 2.6, 0.65,
              fill=NAVY, size=13)
arrow(s, 1.7, 4.0, 1.7, 4.35)
label_in_rect(s, "Bronze Consumer (ECS)\nbuffer → flush every 500", 0.4, 4.35, 2.6, 0.65,
              fill=BLUE, size=12)
arrow(s, 1.7, 5.0, 1.7, 5.35)
label_in_rect(s, "S3 bronze/ (Iceberg)\n+ Glue Data Catalog", 0.4, 5.35, 2.6, 0.65,
              fill=ORANGE, size=13)

# Right: key points
label_in_rect(s, "KEY POINTS", 3.4, 1.3, 9.5, 0.4, fill=NAVY, size=13)

points = [
    ("Kafka keyed by order_id",
     "All messages for an order go to the same partition.\n"
     "Enables ordered processing and efficient partition-level queries."),
    ("At-least-once delivery",
     "enable_auto_commit=False. Offset committed only AFTER successful write.\n"
     "Guarantees no data loss — duplicates absorbed by silver dedup."),
    ("Bronze record = raw + metadata",
     "Keeps all source fields + kafka_topic, kafka_partition, kafka_offset,\n"
     "ingested_at, source_file — full audit trail, never modified."),
    ("CloudBronzeWriter → PyIceberg",
     "In cloud mode: PyIceberg appends PyArrow batches to S3.\n"
     "Table auto-created in Glue on first run (format-version 2, zstd)."),
]
for i, (title, body) in enumerate(points):
    y = 1.9 + i * 1.22
    rect(s, 3.4, y, 9.5, 1.1, fill=WHITE, border=LIGHT_GRAY)
    text_box(s, title, 3.6, y + 0.07, 9.1, 0.33, size=14, bold=True, color=NAVY)
    text_box(s, body,  3.6, y + 0.38, 9.1, 0.65, size=12, color=DARK_GRAY)

add_notes(s, "Le Producer lit le CSV ligne par ligne et envoie chaque SalesRecord vers MSK. La clé Kafka est l'order_id, ce qui garantit que tous les messages d'un même ordre vont sur la même partition. Le Bronze Consumer reçoit ces messages, construit un BronzeRecord avec les métadonnées Kafka, et les écrit en Iceberg via PyIceberg. L'offset n'est commité qu'après l'écriture réussie — at-least-once garanti.")


# ═══════════════════════════════════════════════════════════════════════════ #
# SLIDE 7 — Step 3: Transform Bronze → Silver                                  #
# ═══════════════════════════════════════════════════════════════════════════ #

s = prs.slides.add_slide(blank_layout)
rect(s, 0, 0, W, H, fill=OFF_WHITE)
header_band(s, "Step 3 — Transform Bronze → Silver",
            "Typed · Deduplicated · Partitioned by year/month")

# Bronze schema box
label_in_rect(s, "BRONZE RECORD", 0.4, 1.3, 4.4, 0.4, fill=ORANGE, size=13)
rect(s, 0.4, 1.7, 4.4, 2.9, fill=WHITE, border=LIGHT_GRAY)
bronze_fields = [
    "order_id, region, country, item_type",
    "sales_channel, priority",
    "order_date, ship_date",
    "units_sold, unit_price, unit_cost",
    "total_revenue, total_cost, total_profit",
    "kafka_topic, kafka_partition, kafka_offset",
    "ingested_at, source_file",
]
for i, f in enumerate(bronze_fields):
    text_box(s, f, 0.55, 1.8 + i * 0.37, 4.1, 0.35, size=11, color=DARK_GRAY)

# Arrow
arrow(s, 4.8, 3.1, 5.6, 3.1)
label_in_rect(s, "transform\n_to_silver()", 4.55, 2.85, 1.3, 0.55,
              fill=NAVY, size=10, bold=True)

# Silver schema box
label_in_rect(s, "SILVER RECORD", 5.7, 1.3, 7.2, 0.4, fill=GREEN, size=13)
rect(s, 5.7, 1.7, 7.2, 2.9, fill=WHITE, border=LIGHT_GRAY)
silver_fields = [
    ("All business fields carried over", False, DARK_GRAY),
    ("+ order_year, order_month  [partition keys]", True, NAVY),
    ("+ lead_time_days  [ship - order date]", True, NAVY),
    ("+ margin_pct  [Decimal(6,2), ROUND_HALF_UP]", True, NAVY),
    ("Financials: Decimal(18,4) — no IEEE-754 drift", True, GREEN),
    ("source_kafka_offset  [dedup key]", True, MID_GRAY),
    ("bronze_ingested_at, silver_transformed_at", False, MID_GRAY),
]
for i, (f, bold, color) in enumerate(silver_fields):
    text_box(s, f, 5.85, 1.8 + i * 0.37, 6.9, 0.35, size=11, bold=bold, color=color)

# Bottom row: dedup and partition
label_in_rect(s, "DEDUPLICATION", 0.4, 4.75, 4.4, 0.38, fill=NAVY, size=13)
rect(s, 0.4, 5.13, 4.4, 1.8, fill=WHITE, border=LIGHT_GRAY)
text_box(s,
         "On startup: CloudSilverWriter scans the existing silver table\n"
         "to load all known source_kafka_offset values.\n\n"
         "On merge(): if offset already seen → skip silently.\n"
         "Idempotent replay of any Kafka message — safe for restarts.",
         0.55, 5.2, 4.1, 1.65, size=12, color=DARK_GRAY)

label_in_rect(s, "PARTITIONING", 5.7, 4.75, 7.2, 0.38, fill=NAVY, size=13)
rect(s, 5.7, 5.13, 7.2, 1.8, fill=WHITE, border=LIGHT_GRAY)
text_box(s,
         "Iceberg PartitionSpec: IdentityTransform on (order_year, order_month)\n\n"
         "Athena prunes partitions on date-range queries → cost & speed.\n"
         "silver/year=2020/month=01/  … year=2024/month=12/",
         5.85, 5.2, 6.9, 1.65, size=12, color=DARK_GRAY)

add_notes(s, "La transformation est une pure fonction Python : transform_to_silver(bronze) → silver. Elle calcule les champs dérivés (lead_time_days, margin_pct en Decimal), type tous les montants en Decimal pour éviter les erreurs d'arrondi IEEE-754, et assigne les clés de partition. La dédup se fait sur source_kafka_offset chargé depuis la table Iceberg existante au démarrage. Pas de doublon même si on rejoue les messages Kafka.")


# ═══════════════════════════════════════════════════════════════════════════ #
# SLIDE 8 — Step 4: Serve via Athena                                           #
# ═══════════════════════════════════════════════════════════════════════════ #

s = prs.slides.add_slide(blank_layout)
rect(s, 0, 0, W, H, fill=OFF_WHITE)
header_band(s, "Step 4 — Serve via Athena",
            "Query bronze and silver directly — no crawlers, no ETL jobs")

# Left: query examples
label_in_rect(s, "NAMED QUERIES", 0.4, 1.3, 6.0, 0.4, fill=NAVY, size=13)

queries = [
    ("Preview silver (latest 100 rows)",
     "SELECT * FROM sales_silver\nORDER BY silver_transformed_at DESC LIMIT 100"),
    ("Revenue by country (silver)",
     "SELECT country, SUM(total_revenue) AS revenue\nFROM sales_silver GROUP BY 1 ORDER BY 2 DESC"),
    ("Dedup check",
     "SELECT source_kafka_offset, COUNT(*) c\nFROM sales_silver GROUP BY 1 HAVING c > 1"),
]
for i, (title, sql) in enumerate(queries):
    y = 1.8 + i * 1.65
    rect(s, 0.4, y, 6.0, 1.55, fill=WHITE, border=LIGHT_GRAY)
    text_box(s, title, 0.6, y + 0.05, 5.6, 0.3, size=12, bold=True, color=NAVY)
    rect(s, 0.5, y + 0.38, 5.8, 1.05, fill=DARK_GRAY)
    text_box(s, sql, 0.6, y + 0.42, 5.6, 1.0, size=11, color=OFF_WHITE)

# Right: key points
label_in_rect(s, "WHY ATHENA", 6.8, 1.3, 6.1, 0.4, fill=SKY, size=13)

points = [
    ("No crawler needed",
     "PyIceberg registers tables directly in Glue Data Catalog at write time.\n"
     "Schema always in sync with the data — zero maintenance."),
    ("Iceberg native support",
     "Athena Engine v3 reads Iceberg format natively.\n"
     "Time travel, partition pruning, schema evolution all supported."),
    ("Cost control",
     "Workgroup enforces 10GB scan limit per query (non-prod).\n"
     "Partition pruning on (order_year, order_month) reduces scan size."),
    ("Zero-ops",
     "Serverless — no cluster to manage, pay per TB scanned.\n"
     "Results stored in S3 athena-results/ with 7-day lifecycle."),
]
for i, (title, body) in enumerate(points):
    y = 1.8 + i * 1.22
    rect(s, 6.8, y, 6.1, 1.1, fill=WHITE, border=LIGHT_GRAY)
    text_box(s, title, 7.0, y + 0.07, 5.7, 0.33, size=14, bold=True, color=NAVY)
    text_box(s, body,  7.0, y + 0.38, 5.7, 0.65, size=12, color=DARK_GRAY)

add_notes(s, "Athena Engine v3 peut requêter les tables Iceberg directement depuis Glue — pas besoin de crawler ou de jobs ETL. Les tables sont enregistrées par PyIceberg à l'écriture. J'ai défini 3 named queries préconfigurées dans le workgroup. Le scan limit à 10GB protège les coûts en dev et int. En prod, la limite peut être relevée ou supprimée.")


# ═══════════════════════════════════════════════════════════════════════════ #
# SLIDE 9 — Infrastructure as Code                                             #
# ═══════════════════════════════════════════════════════════════════════════ #

s = prs.slides.add_slide(blank_layout)
rect(s, 0, 0, W, H, fill=OFF_WHITE)
header_band(s, "Infrastructure as Code",
            "Terraform · 7 modules · 4 environments · Zero manual steps")

# Module boxes
modules = [
    ("networking",  "VPC, 3-AZ subnets\nIGW, NAT gateways\nMSK security group",         0.4,  1.5),
    ("s3",          "Data lake bucket\nVersioning, SSE-S3\nLifecycle rules",              0.4,  3.5),
    ("msk",         "MSK cluster\nKafka 3.x\nCloudWatch + Prometheus",                   3.4,  1.5),
    ("glue",        "Glue database\nIAM role for crawlers\nTables via PyIceberg",         3.4,  3.5),
    ("athena",      "Workgroup Engine v3\n10GB scan limit\n3 named queries",              6.4,  1.5),
    ("iam",         "Pipeline execution role\nS3 + MSK + Glue + CW\nLeast privilege",    6.4,  3.5),
    ("scheduler",   "EventBridge Scheduler\nStep Functions SM\nECS Fargate tasks",       9.4,  1.5),
]

for name, desc, x, y in modules:
    color = {"networking": SKY, "s3": ORANGE, "msk": NAVY,
             "glue": GREEN, "athena": SKY, "iam": MID_GRAY,
             "scheduler": BLUE}[name]
    label_in_rect(s, f"module \"{name}\"", x, y, 2.85, 0.38, fill=color, size=12)
    rect(s, x, y + 0.38, 2.85, 1.35, fill=WHITE, border=LIGHT_GRAY)
    text_box(s, desc, x + 0.15, y + 0.45, 2.6, 1.2, size=11, color=DARK_GRAY)

# Environments
label_in_rect(s, "4 ENVIRONMENTS", 0.4, 5.6, 12.5, 0.38, fill=NAVY, size=13)
envs = [
    ("dev.tfvars",  "kafka.t3.small · 1 broker",   "Local iteration"),
    ("int.tfvars",  "kafka.m5.large · 3 brokers",   "Integration tests"),
    ("uat.tfvars",  "kafka.m5.large · 3 brokers",   "Pre-production"),
    ("prod.tfvars", "kafka.m5.2xlarge · 3 brokers",  "Production"),
]
for i, (file, infra, role) in enumerate(envs):
    x = 0.4 + i * 3.15
    rect(s, x, 6.0, 3.0, 1.1, fill=WHITE, border=LIGHT_GRAY)
    text_box(s, file, x + 0.15, 6.05, 2.7, 0.3, size=12, bold=True, color=NAVY)
    text_box(s, infra, x + 0.15, 6.35, 2.7, 0.3, size=11, color=DARK_GRAY)
    text_box(s, role,  x + 0.15, 6.65, 2.7, 0.3, size=11, color=MID_GRAY)

add_notes(s, "7 modules Terraform couvrent toute l'infrastructure. Chaque module a ses propres variables, outputs, et peut être testé indépendamment. Les 4 environnements se distinguent principalement par la taille du cluster MSK. En prod : kafka.m5.2xlarge avec 3 brokers pour la haute disponibilité multi-AZ. En dev : kafka.t3.small avec 1 seul broker pour minimiser les coûts.")


# ═══════════════════════════════════════════════════════════════════════════ #
# SLIDE 10 — Continuous Ingestion                                              #
# ═══════════════════════════════════════════════════════════════════════════ #

s = prs.slides.add_slide(blank_layout)
rect(s, 0, 0, W, H, fill=OFF_WHITE)
header_band(s, "Continuous Ingestion",
            "Designed from day 1 for daily batches — not a one-off script")

# Left: before
label_in_rect(s, "BEFORE  (one-off)", 0.4, 1.3, 5.8, 0.4, fill=MID_GRAY, size=13)
rect(s, 0.4, 1.7, 5.8, 2.3, fill=WHITE, border=LIGHT_GRAY)
before = [
    "• Download runs once manually",
    "• No check: always re-downloads",
    "• No tracking of what was processed",
    "• Services started manually in sequence",
    "• Bronze consumers idle between runs",
]
for i, b in enumerate(before):
    text_box(s, b, 0.6, 1.8 + i * 0.43, 5.4, 0.4, size=13, color=DARK_GRAY)

# Right: after
label_in_rect(s, "AFTER  (continuous)", 6.6, 1.3, 6.3, 0.4, fill=GREEN, size=13)
rect(s, 6.6, 1.7, 6.3, 2.3, fill=WHITE, border=LIGHT_GRAY)
after = [
    "✅  EventBridge fires daily at 09:00 UTC",
    "✅  HTTP HEAD → ETag check before download",
    "✅  IngestionState: extracted → produced",
    "✅  Step Functions chains downloader → producer",
    "✅  Consumers always running, ready for messages",
]
for i, a in enumerate(after):
    text_box(s, a, 6.8, 1.8 + i * 0.43, 5.9, 0.4, size=13, color=DARK_GRAY)

# What doesn't change
label_in_rect(s, "WHAT STAYS THE SAME — zero code changes in services",
              0.4, 4.2, 12.5, 0.4, fill=NAVY, size=13)
rect(s, 0.4, 4.6, 12.5, 2.55, fill=WHITE, border=LIGHT_GRAY)

unchanged = [
    ("Bronze Consumer", "Always running, consumes any message that arrives on MSK"),
    ("Silver Transformer", "Reads bronze, deduplicates, writes silver — unaffected by trigger"),
    ("Athena Serve Layer", "Tables are live — queries work immediately after new data lands"),
    ("Dedup logic", "source_kafka_offset prevents duplicates even on full-dataset replays"),
]
for i, (title, body) in enumerate(unchanged):
    x = 0.6 + (i % 2) * 6.2
    y = 4.7 + (i // 2) * 1.1
    text_box(s, f"✓  {title}", x, y, 5.8, 0.35, size=14, bold=True, color=GREEN)
    text_box(s, body,          x, y + 0.35, 5.8, 0.65, size=12, color=DARK_GRAY)

add_notes(s, "L'architecture est conçue pour la continuité depuis le départ. Le layer Kafka est déjà continu — les consumers tournent en permanence et consomment dès qu'un message arrive. Seul le côté batch (downloader + producer) est déclenché par le scheduler. La clé de l'évolution : on n'a touché qu'au downloader et on a ajouté le scheduler — aucun des services core n'a changé.")


# ═══════════════════════════════════════════════════════════════════════════ #
# SLIDE 11 — Key Design Decisions                                              #
# ═══════════════════════════════════════════════════════════════════════════ #

s = prs.slides.add_slide(blank_layout)
rect(s, 0, 0, W, H, fill=OFF_WHITE)
header_band(s, "Key Design Decisions",
            "Every choice has a clear rationale — and a documented trade-off")

decisions = [
    ("Why Kafka / MSK?",
     "Decouples batch producer from streaming consumers.\n"
     "Replay, exactly-once guarantee, future fan-out to N consumers.",
     "Direct S3 → Lambda: simpler but no replay, no back-pressure."),
    ("Why Iceberg (not Parquet flat)?",
     "ACID writes, schema evolution, time travel, partition pruning.\n"
     "Athena Engine v3 reads it natively — zero glue code needed.",
     "Delta Lake or Hudi: similar features, more complex Athena setup."),
    ("Why Python (not Spark)?",
     "2M rows/day is not big data. Python + Kafka = sufficient.\n"
     "Simple to deploy on ECS Fargate, no cluster cost.",
     "Spark: justified at 100M+ rows or complex SQL transformations."),
    ("Why Pydantic + Decimal?",
     "Type safety at the boundary: CSV strings → typed models.\n"
     "Decimal(18,4) avoids IEEE-754 rounding in financial aggregations.",
     "Dataclasses: no validation. Pandera: more overhead, less portable."),
    ("Why at-least-once + dedup?",
     "Simpler than exactly-once. Kafka exactly-once requires\n"
     "idempotent producer + transactional consumer = more complexity.",
     "If dedup scan at startup is too slow: use a Redis set or DynamoDB."),
]

headers = ["DECISION", "CHOICE & RATIONALE", "ALTERNATIVE CONSIDERED"]
col_widths = [2.5, 6.0, 4.3]
col_x = [0.2, 2.75, 8.8]

# Header row
for j, (h, w, x) in enumerate(zip(headers, col_widths, col_x)):
    label_in_rect(s, h, x, 1.25, w, 0.35, fill=NAVY, size=11)

# Data rows
for i, (dec, choice, alt) in enumerate(decisions):
    y = 1.6 + i * 1.02
    fill = OFF_WHITE if i % 2 == 0 else WHITE
    for j, (w, x) in enumerate(zip(col_widths, col_x)):
        rect(s, x, y, w, 1.0, fill=fill, border=LIGHT_GRAY)
    text_box(s, dec,    col_x[0] + 0.1, y + 0.08, col_widths[0] - 0.15, 0.85, size=12, bold=True, color=NAVY)
    text_box(s, choice, col_x[1] + 0.1, y + 0.08, col_widths[1] - 0.15, 0.85, size=11, color=DARK_GRAY)
    text_box(s, alt,    col_x[2] + 0.1, y + 0.08, col_widths[2] - 0.15, 0.85, size=11, color=MID_GRAY)

add_notes(s, "Ce slide est celui que vous voudrez bien maîtriser pour le Q&A. Les questions les plus probables : pourquoi Kafka et pas juste S3 event trigger ? Réponse : replay, découplage, fan-out futur. Pourquoi Python et pas Spark ? 2M rows/jour c'est trivial pour Python sur Fargate. Si le volume augmentait à 100M+, on switcherait le transformer sur Flink sans changer le DAO.")


# ═══════════════════════════════════════════════════════════════════════════ #
# SLIDE 12 — Code Structure                                                    #
# ═══════════════════════════════════════════════════════════════════════════ #

s = prs.slides.add_slide(blank_layout)
rect(s, 0, 0, W, H, fill=OFF_WHITE)
header_band(s, "Project Structure",
            "Monorepo · Shared package · 4 services · Tests · IaC")

# Directory tree
rect(s, 0.4, 1.25, 5.5, 5.9, fill=DARK_GRAY)
tree = """economics-data-ingestion/
├── packages/
│   └── economics_pipeline/
│       ├── models/        # Pydantic: SalesRecord, BronzeRecord, SilverRecord
│       ├── config/        # PipelineSettings (pydantic-settings, .env.*)
│       ├── dao/           # Local* + Cloud* writers (Iceberg/PyIceberg)
│       ├── kafka/         # SalesProducer, BronzeConsumer
│       ├── transforms/    # transform_to_silver() — pure function
│       ├── exceptions/    # PipelineError hierarchy
│       ├── retry/         # @network_retry @kafka_retry @storage_retry
│       └── ingestion/     # IngestionStateManager (ETag tracking)
├── services/
│   ├── downloader/        # Fetch → archive/ → landing/
│   ├── producer/          # landing/ → MSK topic
│   ├── consumer_bronze/   # MSK → Iceberg bronze
│   ├── transformer_silver/# bronze → Iceberg silver
│   └── scheduler/         # Daily 09h trigger (local dev)
├── tests/
│   ├── unit/              # 95 tests, no external dependencies
│   ├── integration/       # @pytest.mark.integration
│   └── e2e/               # @pytest.mark.e2e
└── infra/
    └── modules/           # networking · s3 · msk · glue · athena · iam · scheduler"""
text_box(s, tree, 0.5, 1.35, 5.3, 5.7, size=9.5, color=OFF_WHITE)

# Right: stats
label_in_rect(s, "BY THE NUMBERS", 6.2, 1.25, 6.7, 0.38, fill=NAVY, size=13)

stats = [
    ("95",   "unit tests passing", "zero external deps"),
    ("7",    "Terraform modules",  "all environments covered"),
    ("4",    "microservices",       "each its own Fargate task"),
    ("3",    "retry policies",     "network / kafka / storage"),
    ("100%", "type-safe",          "Pydantic models, Protocols, Pyright"),
]
for i, (num, label, sub) in enumerate(stats):
    y = 1.7 + i * 1.07
    rect(s, 6.2, y, 6.7, 1.0, fill=WHITE, border=LIGHT_GRAY)
    text_box(s, num,   6.35, y + 0.08, 1.5, 0.55, size=30, bold=True, color=NAVY)
    text_box(s, label, 7.9,  y + 0.15, 4.8, 0.35, size=15, bold=True, color=DARK_GRAY)
    text_box(s, sub,   7.9,  y + 0.5,  4.8, 0.35, size=12, color=MID_GRAY)

add_notes(s, "Le projet est un monorepo avec un package Python partagé installé en mode éditable dans chaque service (pip install -e packages/). Les services sont légers — ils importent depuis economics_pipeline et orchestrent la logique. Les tests sont en 3 tiers : unit (pas d'infra), integration (Kafka ou DAO réel), e2e (pipeline complet). 95 tests unitaires, tous passent sans dépendances externes.")


# ═══════════════════════════════════════════════════════════════════════════ #
# SLIDE 13 — Summary & Q&A                                                     #
# ═══════════════════════════════════════════════════════════════════════════ #

s = prs.slides.add_slide(blank_layout)
rect(s, 0, 0, W, H, fill=NAVY)
rect(s, 0, H - 0.5, W, 0.5, fill=ORANGE)

text_box(s, "Summary", 0.7, 0.4, W - 1.4, 0.7,
         size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
rect(s, 0.7, 1.1, 4.0, 0.04, fill=ORANGE)

# Checklist
checks = [
    "✅  Acquire & Land — daily batch at 09:00, ETag dedup, S3 archive/ + landing/",
    "✅  Stream into Bronze — MSK / Kafka, at-least-once, Iceberg table in Glue",
    "✅  Transform → Silver — typed, deduped, partitioned (year/month)",
    "✅  Serve — Athena Engine v3, named queries, cost-controlled workgroup",
    "✅  IaC — 7 Terraform modules, 4 environments, EventBridge + Step Functions",
]
for i, c in enumerate(checks):
    text_box(s, c, 0.7, 1.3 + i * 0.55, W * 0.6, 0.5,
             size=14, color=WHITE)

# Anticipated Q&A
rect(s, W * 0.63, 1.2, W * 0.36, 5.7, fill=rgb(0x00, 0x1A, 0x40))
text_box(s, "Anticipated Q&A", W * 0.63 + 0.2, 1.3, W * 0.35, 0.4,
         size=14, bold=True, color=ORANGE)
qa = [
    ("Why Kafka and not S3 events?",
     "Replay + decoupling. S3 events are fire-and-forget."),
    ("How does dedup handle scale?",
     "Scan at startup. For 10M+ offsets: Redis or DynamoDB."),
    ("Why Python, not Spark?",
     "2M/day = trivial. Spark justified at 100M+."),
    ("Is the pipeline idempotent?",
     "Yes — ETag check, offset-based filenames, silver dedup."),
]
for i, (q, a) in enumerate(qa):
    y = 1.85 + i * 1.2
    text_box(s, f"Q: {q}", W * 0.63 + 0.2, y, W * 0.35, 0.38,
             size=12, bold=True, color=SKY)
    text_box(s, f"A: {a}", W * 0.63 + 0.2, y + 0.38, W * 0.35, 0.65,
             size=12, color=WHITE)

text_box(s, "Questions?", 0.7, 6.5, 5.0, 0.7,
         size=32, bold=True, color=ORANGE)

add_notes(s, "Pour conclure : les 5 exigences du case study sont toutes couvertes. En 15 minutes on a parcouru l'architecture, les 4 étapes du dataflow, l'IaC, et les choix techniques. Les questions les plus probables sont sur Kafka vs S3 events, et Python vs Spark — vous avez les réponses. Merci!")


# ─── Save ────────────────────────────────────────────────────────────────────

output = "IATA_Economics_Pipeline.pptx"
prs.save(output)
print(f"OK  Saved: {output}  ({prs.slides.__len__()} slides)")
